"""발표용 다중 슬라이드 PPT 생성. 7분 분량, 12 슬라이드.
채점 조건(CI/CD · Docker/Nginx · Web Storage/DBMS)을 각 1슬라이드 이상 깊게 다룸.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ─────────────── 디자인 토큰 ───────────────
NAVY = RGBColor(0x1D, 0x4E, 0xD8)
NAVY_DARK = RGBColor(0x17, 0x3C, 0xA8)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GREY = RGBColor(0x47, 0x55, 0x69)
LIGHT_GREY = RGBColor(0xCB, 0xD5, 0xE1)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
ROSE = RGBColor(0xE1, 0x1D, 0x48)

HANGUL_FONT = "맑은 고딕"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ─────────────── 헬퍼 ───────────────
def set_font(run, *, size=14, bold=False, color=DARK, font=HANGUL_FONT, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(f"{{{NS}}}ea")
    if ea is None:
        ea = etree.SubElement(rPr, f"{{{NS}}}ea")
    ea.set("typeface", font)


def add_text(slide, *, left, top, width, height, text,
             size=14, bold=False, color=DARK, font=HANGUL_FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, font=font, italic=italic)
    return tb


def add_bullets(slide, *, left, top, width, height, items,
                size=14, color=DARK, font=HANGUL_FONT, bullet_color=None, line_gap=4):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_gap)
        if isinstance(item, tuple):
            bullet, rest = item
            r1 = p.add_run()
            r1.text = f"{bullet}  "
            set_font(r1, size=size, bold=True, color=bullet_color or NAVY, font=font)
            r2 = p.add_run()
            r2.text = rest
            set_font(r2, size=size, color=color, font=font)
        else:
            run = p.add_run()
            run.text = f"•  {item}"
            set_font(run, size=size, color=color, font=font)


def add_rect(slide, left, top, width, height, *, fill=LIGHT, line=NAVY, line_pt=1.2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_pt)
    return sh


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, weight=1.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    # 끝에 화살표 머리
    line_elem = conn.line._get_or_add_ln()
    tail = etree.SubElement(line_elem, f"{{{NS}}}tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return conn


def add_header(slide, title, subtitle=None, page_no=None, page_total=None):
    # 상단 바
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill=NAVY, line=None,
             shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, left=Inches(0.6), top=Inches(0.12), width=Inches(10), height=Inches(0.45),
             text=title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, left=Inches(0.6), top=Inches(0.5), width=Inches(10), height=Inches(0.35),
                 text=subtitle, size=11, color=LIGHT_GREY)
    if page_no is not None and page_total is not None:
        add_text(slide, left=Inches(11.7), top=Inches(0.2), width=Inches(1.4), height=Inches(0.4),
                 text=f"{page_no} / {page_total}", size=11, color=LIGHT_GREY,
                 align=PP_ALIGN.RIGHT)


def add_footer(slide, page_no=None, total=None):
    add_text(slide, left=Inches(0.6), top=Inches(7.1), width=Inches(8), height=Inches(0.3),
             text="PlanIt — 2026-1 Web Programming Final Project · 202100580 이정",
             size=9, color=GREY)
    if page_no and total:
        add_text(slide, left=Inches(11.5), top=Inches(7.1), width=Inches(1.4), height=Inches(0.3),
                 text=f"{page_no} / {total}", size=9, color=GREY, align=PP_ALIGN.RIGHT)


# ─────────────── 개별 슬라이드 ───────────────
def slide_cover(pres):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    # 배경
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY, line=None, shape=MSO_SHAPE.RECTANGLE)
    # 부제 작게
    add_text(s, left=Inches(0.8), top=Inches(1.4), width=Inches(11), height=Inches(0.4),
             text="2026-1 Web Programming · Final Project", size=14, color=LIGHT_GREY)
    # 큰 제목
    add_text(s, left=Inches(0.8), top=Inches(2.0), width=Inches(11), height=Inches(1.6),
             text="PlanIt", size=96, bold=True, color=WHITE)
    add_text(s, left=Inches(0.8), top=Inches(3.7), width=Inches(11), height=Inches(0.6),
             text="할 일 + 일정 통합 관리 웹 서비스", size=28, color=AMBER)
    add_text(s, left=Inches(0.8), top=Inches(4.5), width=Inches(11), height=Inches(0.5),
             text="대학생을 위한 단일 인터페이스 · React + Express + PostgreSQL · Docker · Nginx · CI/CD",
             size=14, color=LIGHT_GREY)
    # 학생 정보
    add_text(s, left=Inches(0.8), top=Inches(6.2), width=Inches(11), height=Inches(0.4),
             text="컴퓨터공학부 · 202100580 · 이정", size=16, bold=True, color=WHITE)


def slide_problem(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "왜 만들었는가", "대학생의 두 가지 일정 패턴", page, total)

    # 좌: 문제, 우: 솔루션
    add_rect(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.2), fill=LIGHT, line=ROSE)
    add_text(s, left=Inches(0.85), top=Inches(1.45), width=Inches(5.6), height=Inches(0.4),
             text="❗ 문제", size=14, bold=True, color=ROSE)
    add_text(s, left=Inches(0.85), top=Inches(1.9), width=Inches(5.6), height=Inches(0.8),
             text="대학생은 매일 두 종류의 일정을 다룬다", size=18, bold=True, color=DARK)
    add_bullets(s, left=Inches(0.85), top=Inches(2.85), width=Inches(5.6), height=Inches(3.5),
                items=[
                    ("Event", "강의 · 미팅 · 시험 → 시각이 고정된 일정"),
                    ("Task", "과제 · 보고서 · 학습 → 마감일이 있는 작업"),
                    ("현실", "캘린더 앱과 투두 앱이 분리 → 두 앱을 오가야 함"),
                    ("결과", "맥락 전환 비용, 누락 위험"),
                ],
                size=13, line_gap=8)

    add_rect(s, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.2), fill=LIGHT, line=EMERALD)
    add_text(s, left=Inches(7.15), top=Inches(1.45), width=Inches(5.6), height=Inches(0.4),
             text="✔ 해결", size=14, bold=True, color=EMERALD)
    add_text(s, left=Inches(7.15), top=Inches(1.9), width=Inches(5.6), height=Inches(0.8),
             text="둘을 한 화면에 모은다", size=18, bold=True, color=DARK)
    add_bullets(s, left=Inches(7.15), top=Inches(2.85), width=Inches(5.6), height=Inches(3.5),
                items=[
                    "Task 와 Event 를 동일한 데이터 모델로 추상화",
                    "대시보드 = 오늘 마감 할 일 + 오늘 일정",
                    "통합 캘린더 = 월/주/일/아젠다 한 화면",
                    "대학생 → 일반 사용자로 확장 가능한 단순 구조",
                ],
                size=13, line_gap=8)

    add_footer(s, page, total)


def slide_overview(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "서비스 한눈에", "핵심 기능 6가지", page, total)

    items = [
        ("🔐", "인증", "이메일/비밀번호 + bcrypt + JWT 7일"),
        ("✅", "할 일 (Task)", "제목 · 마감일 · 우선순위 · 완료 · 검색/필터"),
        ("📅", "일정 (Event)", "제목 · 장소 · 시작/종료"),
        ("🗓️", "통합 캘린더", "월/주/일/아젠다 + 한국어 + Task/Event 색 분리"),
        ("🏠", "대시보드", "오늘 마감 할 일 + 오늘 일정"),
        ("🔔", "마감 알림 · 다크 모드", "Notification API + localStorage 영구"),
    ]
    cols, rows = 3, 2
    cell_w = Inches(4.05)
    cell_h = Inches(2.4)
    gap = Inches(0.15)
    x0 = Inches(0.6)
    y0 = Inches(1.2)
    for idx, (emoji, name, desc) in enumerate(items):
        r, c = divmod(idx, cols)
        x = x0 + (cell_w + gap) * c
        y = y0 + (cell_h + gap) * r
        add_rect(s, x, y, cell_w, cell_h, fill=WHITE, line=NAVY)
        add_text(s, left=x + Inches(0.2), top=y + Inches(0.15), width=Inches(1.0), height=Inches(0.7),
                 text=emoji, size=28, color=NAVY)
        add_text(s, left=x + Inches(0.2), top=y + Inches(1.0), width=cell_w - Inches(0.4), height=Inches(0.5),
                 text=name, size=16, bold=True, color=DARK)
        add_text(s, left=x + Inches(0.2), top=y + Inches(1.5), width=cell_w - Inches(0.4), height=Inches(0.8),
                 text=desc, size=11, color=GREY)

    add_footer(s, page, total)


def slide_architecture(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "전체 아키텍처", "Browser → Vercel → Render(Docker) → Neon Postgres", page, total)

    # 박스 4개: Browser, Vercel/Nginx, Render API, Neon DB
    boxes = [
        ("Browser",         "React SPA\nlocalStorage / sessionStorage / IndexedDB", AMBER),
        ("Vercel CDN\n(Nginx-style)", "정적 자원 + SPA fallback\nimmutable cache + gzip", NAVY),
        ("Render Web Service\n(Docker)", "Express + Prisma\nJWT · bcrypt · zod", EMERALD),
        ("Neon Postgres", "User · Task · Event\n1:N 관계", ROSE),
    ]
    n = len(boxes)
    bw = Inches(2.7)
    bh = Inches(2.0)
    gap_x = Inches(0.35)
    total_w = bw * n + gap_x * (n - 1)
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.6)
    centers = []
    for i, (title, desc, color) in enumerate(boxes):
        x = start_x + (bw + gap_x) * i
        add_rect(s, x, y, bw, bh, fill=WHITE, line=color, line_pt=2)
        add_text(s, left=x + Inches(0.15), top=y + Inches(0.2), width=bw - Inches(0.3), height=Inches(0.85),
                 text=title, size=14, bold=True, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left=x + Inches(0.15), top=y + Inches(1.1), width=bw - Inches(0.3), height=Inches(0.85),
                 text=desc, size=10, color=GREY, align=PP_ALIGN.CENTER)
        centers.append((x + bw, y + bh // 2))

    # 화살표
    for i in range(n - 1):
        x_end, y_mid = centers[i]
        next_x = start_x + (bw + gap_x) * (i + 1)
        add_arrow(s, x_end, y_mid, next_x, y_mid, color=NAVY, weight=2)

    # 하단: 요청 흐름 설명
    add_text(s, left=Inches(0.6), top=Inches(5.2), width=Inches(12.1), height=Inches(0.4),
             text="요청 흐름", size=14, bold=True, color=NAVY)
    add_bullets(s, left=Inches(0.6), top=Inches(5.6), width=Inches(12.1), height=Inches(1.4),
                items=[
                    ("①", "사용자는 Vercel 도메인으로 접속 → React 정적 자원이 즉시 응답 (Nginx 식 gzip + immutable cache)"),
                    ("②", "/api/* 요청은 axios 가 Render 백엔드로 직접 (build 시 VITE_API_BASE_URL inline)"),
                    ("③", "Express 는 JWT 검증 → Prisma → Neon Postgres SQL"),
                    ("④", "응답 데이터는 IndexedDB 에 캐시 → 오프라인/네트워크 장애 시 폴백"),
                ],
                size=12, line_gap=4)
    add_footer(s, page, total)


def slide_stack(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "기술 스택", "수업에서 다룬 요소 + 표준 브라우저 API", page, total)

    rows = [
        ("Frontend",  "React · Vite · Tailwind · react-router-dom · react-big-calendar · axios · date-fns"),
        ("Backend",   "Node 20 · Express · Prisma · JWT · bcrypt · zod · cors · morgan"),
        ("Database",  "PostgreSQL 16 (Neon) · Prisma Migrate"),
        ("Container", "Docker (multi-stage) · docker compose (db/api/web 3 컨테이너)"),
        ("Reverse Proxy", "Nginx (gzip · immutable cache · SPA fallback · /api 리버스 프록시)"),
        ("Web Storage", "localStorage · sessionStorage · IndexedDB · Notification API"),
        ("CI/CD",     "GitHub Actions (ci.yml + deploy.yml) · Vercel · Render autoDeploy"),
    ]
    y0 = Inches(1.2)
    row_h = Inches(0.78)
    for i, (k, v) in enumerate(rows):
        y = y0 + row_h * i
        # zebra
        fill = LIGHT if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.6), y, Inches(12.1), row_h, fill=fill, line=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, left=Inches(0.85), top=y + Inches(0.18), width=Inches(2.6), height=row_h,
                 text=k, size=14, bold=True, color=NAVY)
        add_text(s, left=Inches(3.5), top=y + Inches(0.18), width=Inches(9.0), height=row_h,
                 text=v, size=13, color=DARK)
    add_footer(s, page, total)


def slide_dbms(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "DBMS 활용 — PostgreSQL + Prisma + Neon", "조건 ③ — 영속 데이터의 1:N 관계 모델", page, total)

    # 좌: 스키마
    add_rect(s, Inches(0.6), Inches(1.2), Inches(7.2), Inches(5.4), fill=LIGHT, line=NAVY)
    add_text(s, left=Inches(0.85), top=Inches(1.35), width=Inches(7.0), height=Inches(0.4),
             text="📐 도메인 모델 (Prisma schema)", size=14, bold=True, color=NAVY)
    schema_text = (
        "model User  { id  email(unique)  password  name?  ... }\n"
        "model Task  { id  userId  title  description?  dueAt?\n"
        "              priority(LOW|MEDIUM|HIGH)  completed  ... }\n"
        "model Event { id  userId  title  location?  startAt  endAt  ... }\n"
        "\n"
        "User 1 ──< Task   (마감일이 있는 작업)\n"
        "User 1 ──< Event  (시각이 고정된 일정)\n"
        "\n"
        "@@index([userId, dueAt])           ─ 대시보드 쿼리 가속\n"
        "@@index([userId, startAt])         ─ 캘린더 범위 쿼리\n"
        "onDelete: Cascade                  ─ 계정 삭제 시 정리"
    )
    tb = s.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(7.0), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(schema_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        set_font(run, size=12, color=DARK, font="Consolas")

    # 우: 운영 포인트
    add_rect(s, Inches(8.0), Inches(1.2), Inches(4.7), Inches(5.4), fill=WHITE, line=NAVY)
    add_text(s, left=Inches(8.2), top=Inches(1.35), width=Inches(4.4), height=Inches(0.4),
             text="🚀 운영 포인트", size=14, bold=True, color=NAVY)
    add_bullets(s, left=Inches(8.2), top=Inches(1.85), width=Inches(4.4), height=Inches(4.7),
                items=[
                    ("Prisma", "스키마 우선 · 타입 안전 · migrate deploy 로 운영 일치"),
                    ("Neon",   "Serverless Postgres · 무료 플랜 · Connection Pooling"),
                    ("SSL",    "sslmode=require + channel_binding=require"),
                    ("Index",  "사용자 + 마감/시작 시각 복합 인덱스"),
                    ("ORM 가드", "zod 유효성 + Prisma 타입 = SQL 인젝션 차단"),
                    ("Migration", "prisma/migrations/20260607070117_init"),
                ],
                size=12, line_gap=6)

    add_footer(s, page, total)


def slide_web_storage(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "Web Storage 활용 — 4계층 저장소 분담", "조건 ③ — 데이터 수명 · 범위에 맞춘 분리", page, total)

    rows = [
        ("PostgreSQL",     "영속 · 다기기 일관성",          "User · Task · Event 본체",                        NAVY),
        ("localStorage",   "영구 · 같은 브라우저",           "JWT (planit.token) · 테마 (planit.theme) · 알림 발송 기록", EMERALD),
        ("sessionStorage", "탭 수명 · 새로고침 보존",         "할 일 검색어/우선순위 · 캘린더 마지막 본 월",         AMBER),
        ("IndexedDB",      "비동기 객체 저장 · 오프라인 캐시", "tasks/events 캐시 → 네트워크 실패 시 폴백",         ROSE),
    ]
    # 표 헤더
    y0 = Inches(1.25)
    row_h = Inches(1.05)
    cols = [
        ("저장소",   Inches(0.6),   Inches(2.4)),
        ("성격",     Inches(3.05),  Inches(3.0)),
        ("실제 활용 위치", Inches(6.15), Inches(6.55)),
    ]
    # 헤더
    add_rect(s, Inches(0.6), y0 - Inches(0.55), Inches(12.1), Inches(0.55),
             fill=NAVY, line=None, shape=MSO_SHAPE.RECTANGLE)
    for label, x, w in cols:
        add_text(s, left=x + Inches(0.15), top=y0 - Inches(0.5), width=w, height=Inches(0.45),
                 text=label, size=12, bold=True, color=WHITE)
    # 본문
    for i, (k, kind, usage, color) in enumerate(rows):
        y = y0 + row_h * i
        fill = LIGHT if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.6), y, Inches(12.1), row_h, fill=fill, line=None,
                 shape=MSO_SHAPE.RECTANGLE)
        # 좌측 색 막대
        add_rect(s, Inches(0.6), y, Inches(0.1), row_h, fill=color, line=None,
                 shape=MSO_SHAPE.RECTANGLE)
        # 컬럼 텍스트
        add_text(s, left=cols[0][1] + Inches(0.15), top=y + Inches(0.2), width=cols[0][2], height=row_h,
                 text=k, size=14, bold=True, color=color)
        add_text(s, left=cols[1][1] + Inches(0.15), top=y + Inches(0.2), width=cols[1][2], height=row_h,
                 text=kind, size=12, color=DARK)
        add_text(s, left=cols[2][1] + Inches(0.15), top=y + Inches(0.2), width=cols[2][2], height=row_h,
                 text=usage, size=12, color=DARK)

    # 코드 발췌
    add_rect(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.25), fill=DARK, line=None)
    code = "client/src/api/client.js  ─  JWT 자동 첨부\n" \
           "client/src/lib/idbCache.js ─  IndexedDB 캐시 read/write\n" \
           "client/src/pages/Tasks.jsx ─  sessionStorage 로 검색어 유지\n" \
           "client/src/hooks/useDueNotifier.js ─  Notification API + localStorage 중복 방지"
    tb = s.shapes.add_textbox(Inches(0.85), Inches(5.8), Inches(11.6), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        set_font(run, size=11, color=LIGHT, font="Consolas")
    add_footer(s, page, total)


def slide_docker(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "Docker 컨테이너 구성", "조건 ② — 한 번의 compose up 으로 운영 환경 재현", page, total)

    # 좌측: 다이어그램
    add_rect(s, Inches(0.6), Inches(1.2), Inches(6.6), Inches(5.4), fill=LIGHT, line=NAVY)
    add_text(s, left=Inches(0.85), top=Inches(1.35), width=Inches(6.2), height=Inches(0.4),
             text="🐳 docker-compose 토폴로지", size=14, bold=True, color=NAVY)

    cx = Inches(0.6) + Inches(3.3)
    # web
    web_y = Inches(1.95)
    add_rect(s, cx - Inches(2.4), web_y, Inches(4.8), Inches(0.8), fill=WHITE, line=AMBER, line_pt=2)
    add_text(s, left=cx - Inches(2.4), top=web_y + Inches(0.15), width=Inches(4.8), height=Inches(0.5),
             text="web — Nginx (planit-web)\n:8080 → 80 · 정적 + /api 프록시 + gzip",
             size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    # api
    api_y = Inches(3.25)
    add_rect(s, cx - Inches(2.4), api_y, Inches(4.8), Inches(0.8), fill=WHITE, line=EMERALD, line_pt=2)
    add_text(s, left=cx - Inches(2.4), top=api_y + Inches(0.15), width=Inches(4.8), height=Inches(0.5),
             text="api — Express + Prisma (planit-api)\n:4000 · migrate deploy → node src/index.js",
             size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    # db
    db_y = Inches(4.55)
    add_rect(s, cx - Inches(2.4), db_y, Inches(4.8), Inches(0.8), fill=WHITE, line=ROSE, line_pt=2)
    add_text(s, left=cx - Inches(2.4), top=db_y + Inches(0.15), width=Inches(4.8), height=Inches(0.5),
             text="db — postgres:16-alpine\n:5434 → 5432 · planit-db 볼륨 영속",
             size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # 화살표
    add_arrow(s, cx, web_y + Inches(0.8), cx, api_y, color=NAVY, weight=1.6)
    add_arrow(s, cx, api_y + Inches(0.8), cx, db_y, color=NAVY, weight=1.6)

    add_text(s, left=Inches(0.85), top=Inches(5.55), width=Inches(6.2), height=Inches(0.9),
             text="depends_on healthcheck → db 가 ready 일 때만 api 시작\nnginx.conf 는 volume mount → 재빌드 없이 라우팅 수정 가능",
             size=10, color=GREY, italic=True)

    # 우측: 기술 포인트
    add_rect(s, Inches(7.4), Inches(1.2), Inches(5.3), Inches(5.4), fill=WHITE, line=NAVY)
    add_text(s, left=Inches(7.6), top=Inches(1.35), width=Inches(5.0), height=Inches(0.4),
             text="🔧 기술 포인트", size=14, bold=True, color=NAVY)
    add_bullets(s, left=Inches(7.6), top=Inches(1.85), width=Inches(5.0), height=Inches(4.7),
                items=[
                    ("Multi-stage build", "deps → build (prisma generate) → runtime · 이미지 슬림화"),
                    ("node:20-slim", "Prisma + Alpine OpenSSL 이슈 회피 → 안정적 Debian 베이스"),
                    ("binaryTargets", "debian-openssl-3.0.x 명시로 production 런타임 호환"),
                    ("Health gate", "pg_isready 헬스체크 통과 후 api 부팅"),
                    ("Volume",    "planit-db 명명 볼륨으로 데이터 영속"),
                    ("Port isolation", "호스트 5434 매핑 → 호스트 Postgres 충돌 회피"),
                ],
                size=11, line_gap=4)

    add_footer(s, page, total)


def slide_nginx(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "Nginx 리버스 프록시", "조건 ② — 정적 자원 + API 를 같은 출처에서 서비스", page, total)

    # 좌: 핵심 설정
    add_rect(s, Inches(0.6), Inches(1.2), Inches(7.5), Inches(5.4), fill=DARK, line=None)
    add_text(s, left=Inches(0.85), top=Inches(1.35), width=Inches(7.0), height=Inches(0.4),
             text="docker/nginx.conf 핵심", size=14, bold=True, color=AMBER)
    code = """server {
    listen 80;
    gzip on;
    gzip_types text/css application/javascript application/json image/svg+xml;

    # API 리버스 프록시 (같은 출처 → CORS 회피)
    location /api/ {
        proxy_pass http://api:4000;
        proxy_set_header X-Real-IP $remote_addr;
        proxy_set_header X-Forwarded-For $proxy_add_x_forwarded_for;
    }

    # 정적 자원: 1년 immutable 캐시
    location /assets/ {
        add_header Cache-Control "public, max-age=31536000, immutable";
        try_files $uri =404;
    }

    # SPA fallback (react-router 새로고침 보장)
    location / {
        try_files $uri $uri/ /index.html;
    }
}"""
    tb = s.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(7.0), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        set_font(run, size=10, color=LIGHT, font="Consolas")

    # 우: 효과
    add_rect(s, Inches(8.3), Inches(1.2), Inches(4.4), Inches(5.4), fill=LIGHT, line=NAVY)
    add_text(s, left=Inches(8.5), top=Inches(1.35), width=Inches(4.1), height=Inches(0.4),
             text="🎯 얻은 효과", size=14, bold=True, color=NAVY)
    add_bullets(s, left=Inches(8.5), top=Inches(1.85), width=Inches(4.1), height=Inches(4.7),
                items=[
                    ("같은 Origin", "CORS preflight 불필요"),
                    ("gzip", "JS 429KB → 138KB (전송량 -68%)"),
                    ("immutable", "재방문 시 자산 0 byte 다운로드"),
                    ("SPA fallback", "/calendar 직접 새로고침해도 200"),
                    ("X-Forwarded", "백엔드에서 원본 IP 식별 가능"),
                    ("HTTPS 전담", "운영 환경에서 TLS 처리 Nginx 단계 집중"),
                ],
                size=12, line_gap=6)

    add_footer(s, page, total)


def slide_cicd(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "CI/CD 자동 배포 파이프라인", "조건 ① — git push → 자동 빌드 · 자동 배포", page, total)

    # 흐름도: PR → CI / main push → CI + auto deploy
    add_text(s, left=Inches(0.6), top=Inches(1.15), width=Inches(12.1), height=Inches(0.4),
             text="git push origin main", size=14, bold=True, color=NAVY)

    # 박스 3개 흐름
    flow = [
        ("GitHub Actions\nci.yml",       "server / client / docker compose\n빌드 + node --check + dist 아티팩트", NAVY),
        ("Vercel\nGitHub 연동",          "client/ 자동 감지\n빌드 + Edge 배포 + alias 유지",                   AMBER),
        ("Render\nautoDeploy: yes",      "server/Dockerfile 자동 감지\n빌드 + migrate deploy + 컨테이너 부팅", EMERALD),
    ]
    bw = Inches(3.9)
    bh = Inches(2.1)
    y = Inches(1.7)
    x0 = Inches(0.6)
    gap = Inches(0.2)
    centers = []
    for i, (title, desc, color) in enumerate(flow):
        x = x0 + (bw + gap) * i
        add_rect(s, x, y, bw, bh, fill=WHITE, line=color, line_pt=2)
        add_text(s, left=x + Inches(0.2), top=y + Inches(0.2), width=bw - Inches(0.4), height=Inches(0.8),
                 text=title, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, left=x + Inches(0.2), top=y + Inches(1.05), width=bw - Inches(0.4), height=Inches(0.95),
                 text=desc, size=11, color=DARK, align=PP_ALIGN.CENTER)
        centers.append((x + bw // 2, y))

    # 위에서 화살표
    for cx, cy in centers:
        add_arrow(s, cx, Inches(1.45), cx, cy, color=NAVY, weight=1.4)

    # 하단: 검증/혜택
    add_rect(s, Inches(0.6), Inches(4.2), Inches(6.0), Inches(2.6), fill=LIGHT, line=NAVY)
    add_text(s, left=Inches(0.85), top=Inches(4.3), width=Inches(5.6), height=Inches(0.4),
             text="📋 ci.yml 작업 (PR & main push 시)", size=13, bold=True, color=NAVY)
    add_bullets(s, left=Inches(0.85), top=Inches(4.75), width=Inches(5.6), height=Inches(2.0),
                items=[
                    ("server",   "npm ci + prisma generate + node --check"),
                    ("client",   "npm ci + vite build + dist 아티팩트 업로드"),
                    ("docker",   "docker compose build 전체 검증"),
                    ("CodeRabbit", "PR 자동 코드 리뷰 (선택)"),
                ],
                size=11, line_gap=4)

    add_rect(s, Inches(6.8), Inches(4.2), Inches(5.9), Inches(2.6), fill=LIGHT, line=EMERALD)
    add_text(s, left=Inches(7.05), top=Inches(4.3), width=Inches(5.5), height=Inches(0.4),
             text="🚀 자동 배포 흐름", size=13, bold=True, color=EMERALD)
    add_bullets(s, left=Inches(7.05), top=Inches(4.75), width=Inches(5.5), height=Inches(2.0),
                items=[
                    "main 머지 → Vercel/Render 가 GitHub Webhook 으로 감지",
                    "프런트: Vercel Edge 글로벌 배포 (≈ 30초)",
                    "백엔드: Render 가 Docker 이미지 빌드 + 마이그레이션 자동 적용",
                    "결과: 4개 PR 머지 모두 무중단으로 운영 반영",
                ],
                size=11, line_gap=4)

    add_footer(s, page, total)


def slide_security(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "보안 · 안정성 · 사용자 경험", "기술적 마무리 디테일", page, total)

    items = [
        ("🔒 bcrypt",   "salt round 10 → 평문 비밀번호 영속 저장 없음"),
        ("🪪 JWT",      "HS256 · Bearer · 7일 만료 · localStorage + axios 인터셉터 자동 첨부"),
        ("🛡 zod",      "모든 입력 서버단 유효성 검증 (회원가입/CRUD 페이로드)"),
        ("🚷 401 핸들링","axios 응답 인터셉터 → localStorage 제거 + 로그인 화면 이동"),
        ("🌐 CORS",     "Render env CORS_ORIGIN 화이트리스트 → Vercel 도메인만 허용"),
        ("📴 오프라인",  "IndexedDB 캐시 → 네트워크 실패 시 마지막 데이터 유지"),
        ("🌓 다크 모드", "OS prefers-color-scheme 감지 + 사용자 선택 영구 저장"),
        ("🔔 마감 알림", "1분 주기 폴링 · 1시간 이내 마감 · 중복 알림 방지"),
    ]
    cols = 2
    cell_w = Inches(6.1)
    cell_h = Inches(1.3)
    gap = Inches(0.15)
    x0 = Inches(0.6)
    y0 = Inches(1.25)
    for i, (title, desc) in enumerate(items):
        r, c = divmod(i, cols)
        x = x0 + (cell_w + gap) * c
        y = y0 + (cell_h + gap) * r
        add_rect(s, x, y, cell_w, cell_h, fill=LIGHT, line=NAVY)
        add_text(s, left=x + Inches(0.2), top=y + Inches(0.15), width=cell_w - Inches(0.4), height=Inches(0.5),
                 text=title, size=14, bold=True, color=NAVY)
        add_text(s, left=x + Inches(0.2), top=y + Inches(0.65), width=cell_w - Inches(0.4), height=Inches(0.6),
                 text=desc, size=11, color=DARK)

    add_footer(s, page, total)


def slide_demo(pres, page, total):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "라이브 시연", "이 슬라이드에서 실제 사이트로 전환", page, total)

    # 좌측 큰 안내
    add_rect(s, Inches(0.6), Inches(1.25), Inches(6.5), Inches(5.5), fill=NAVY, line=None)
    add_text(s, left=Inches(0.9), top=Inches(1.5), width=Inches(6.0), height=Inches(0.45),
             text="시연 URL", size=14, bold=True, color=AMBER)
    add_text(s, left=Inches(0.9), top=Inches(2.0), width=Inches(6.0), height=Inches(0.6),
             text="https://client-eight-ivory-60.vercel.app", size=18, bold=True, color=WHITE)
    add_text(s, left=Inches(0.9), top=Inches(3.0), width=Inches(6.0), height=Inches(0.45),
             text="GitHub", size=14, bold=True, color=AMBER)
    add_text(s, left=Inches(0.9), top=Inches(3.45), width=Inches(6.0), height=Inches(0.6),
             text="github.com/Pionia5375/2026-1-web-programming", size=15, bold=True, color=WHITE)
    add_text(s, left=Inches(0.9), top=Inches(4.5), width=Inches(6.0), height=Inches(0.45),
             text="API", size=14, bold=True, color=AMBER)
    add_text(s, left=Inches(0.9), top=Inches(4.95), width=Inches(6.0), height=Inches(0.6),
             text="planit-api-i1th.onrender.com", size=15, bold=True, color=WHITE)
    add_text(s, left=Inches(0.9), top=Inches(6.0), width=Inches(6.0), height=Inches(0.4),
             text="(Render Free — 첫 호출 cold start 30~50초)", size=10, color=LIGHT_GREY, italic=True)

    # 우측 시연 순서
    add_rect(s, Inches(7.3), Inches(1.25), Inches(5.4), Inches(5.5), fill=LIGHT, line=NAVY)
    add_text(s, left=Inches(7.5), top=Inches(1.4), width=Inches(5.0), height=Inches(0.45),
             text="🎬 시연 순서 (2분 15초)", size=14, bold=True, color=NAVY)
    add_bullets(s, left=Inches(7.5), top=Inches(1.95), width=Inches(5.0), height=Inches(4.7),
                items=[
                    ("로그인",    "JWT 발급 → 대시보드 진입"),
                    ("할 일",     "추가 → 우선순위 색 뱃지 → 검색"),
                    ("일정",      "장소/시간 입력 → zod 유효성"),
                    ("캘린더",    "월/주/일/아젠다 + 한국어"),
                    ("다크 모드", "토글 → 새로고침해도 유지"),
                    ("마감 알림", "미리 등록한 할 일 → Notification"),
                    ("IndexedDB", "DevTools 로 planit-cache 확인"),
                ],
                size=12, line_gap=6)
    add_footer(s, page, total)


def slide_closing(pres):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    # 배경
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY, line=None, shape=MSO_SHAPE.RECTANGLE)

    add_text(s, left=Inches(0.8), top=Inches(1.5), width=Inches(12), height=Inches(0.6),
             text="감사합니다.", size=48, bold=True, color=WHITE)
    add_text(s, left=Inches(0.8), top=Inches(2.4), width=Inches(12), height=Inches(0.5),
             text="PlanIt — 14주차 설계서를 그대로 구현한 풀스택 웹 서비스", size=18, color=AMBER)

    # 박스 두 줄: URL
    add_rect(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(1.0), fill=NAVY_DARK, line=AMBER)
    add_text(s, left=Inches(1.1), top=Inches(3.55), width=Inches(2.5), height=Inches(0.6),
             text="🌐 서비스", size=14, bold=True, color=AMBER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left=Inches(3.6), top=Inches(3.55), width=Inches(9.0), height=Inches(0.7),
             text="https://client-eight-ivory-60.vercel.app", size=18, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.0), fill=NAVY_DARK, line=AMBER)
    add_text(s, left=Inches(1.1), top=Inches(4.75), width=Inches(2.5), height=Inches(0.6),
             text="📦 GitHub", size=14, bold=True, color=AMBER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left=Inches(3.6), top=Inches(4.75), width=Inches(9.0), height=Inches(0.7),
             text="github.com/Pionia5375/2026-1-web-programming", size=18, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, left=Inches(0.8), top=Inches(6.4), width=Inches(12), height=Inches(0.5),
             text="컴퓨터공학부 · 202100580 · 이정 · 2026-1 Web Programming", size=14, color=LIGHT_GREY)


# ─────────────── 빌드 진입점 ───────────────
def main():
    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H

    total = 13
    slide_cover(pres)
    slide_problem(pres, 2, total)
    slide_overview(pres, 3, total)
    slide_architecture(pres, 4, total)
    slide_stack(pres, 5, total)
    slide_dbms(pres, 6, total)
    slide_web_storage(pres, 7, total)
    slide_docker(pres, 8, total)
    slide_nginx(pres, 9, total)
    slide_cicd(pres, 10, total)
    slide_security(pres, 11, total)
    slide_demo(pres, 12, total)
    slide_closing(pres)

    out = "/Users/ijeong/Desktop/웹 프로그래밍/2026-1-web-programming/docs/presentation_deck.pptx"
    pres.save(out)
    print(f"saved: {out} · slides: {len(pres.slides)}")


if __name__ == "__main__":
    main()
