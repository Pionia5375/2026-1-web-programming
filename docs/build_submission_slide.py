"""한 페이지 제출용 PPT 생성. 채점에 필요한 GitHub Repo / 서비스 URL 포함."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


HANGUL_FONT = "맑은 고딕"        # Windows 기본
HANGUL_FONT_FALLBACK = "Apple SD Gothic Neo"  # mac 폴백 (East Asia 자동)

NAVY = RGBColor(0x1D, 0x4E, 0xD8)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GREY = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)


def set_font(run, *, size=14, bold=False, color=DARK, font=HANGUL_FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # East Asia font 별도 설정
    rPr = run._r.get_or_add_rPr()
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ea = rPr.find(f'{{{ns}}}ea')
    if ea is None:
        ea = etree.SubElement(rPr, f'{{{ns}}}ea')
    ea.set('typeface', font)


def add_text(slide, *, left, top, width, height, text, size=14, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=HANGUL_FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, font=font)
    return tb


def add_bullet_block(slide, *, left, top, width, height, items, size=12, color=DARK, font=HANGUL_FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = f"• {item}"
        set_font(run, size=size, color=color, font=font)


def main():
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    blank = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank)

    # 배경: 좌측 navy 사이드바
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(3.6), Inches(7.5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = NAVY
    sidebar.line.fill.background()

    # 사이드바: 프로젝트 이름
    add_text(slide, left=Inches(0.4), top=Inches(0.6), width=Inches(3.0), height=Inches(0.7),
             text="PlanIt", size=44, bold=True, color=WHITE)
    add_text(slide, left=Inches(0.4), top=Inches(1.45), width=Inches(3.0), height=Inches(0.4),
             text="할 일 + 일정 통합 웹 서비스", size=13, color=LIGHT)

    # 사이드바: 학생 정보 (하단)
    add_text(slide, left=Inches(0.4), top=Inches(5.5), width=Inches(3.0), height=Inches(0.3),
             text="컴퓨터공학부", size=10, color=LIGHT)
    add_text(slide, left=Inches(0.4), top=Inches(5.85), width=Inches(3.0), height=Inches(0.5),
             text="202100580  이정", size=18, bold=True, color=WHITE)
    add_text(slide, left=Inches(0.4), top=Inches(6.45), width=Inches(3.0), height=Inches(0.3),
             text="2026-1 Web Programming", size=10, color=LIGHT)
    add_text(slide, left=Inches(0.4), top=Inches(6.75), width=Inches(3.0), height=Inches(0.3),
             text="기말 프로젝트", size=10, color=LIGHT)

    # 본문 시작 X
    BX = Inches(4.0)
    BW = Inches(8.9)

    # 본문: 헤드라인
    add_text(slide, left=BX, top=Inches(0.6), width=BW, height=Inches(0.5),
             text="대학생을 위한 통합 일정 관리 서비스", size=22, bold=True, color=DARK)
    add_text(slide, left=BX, top=Inches(1.15), width=BW, height=Inches(0.4),
             text="강의·미팅·시험(Event) 과 과제·보고서(Task) 를 한 화면에서.", size=12, color=GREY)

    # 박스 1: 배포 정보
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, BX, Inches(1.75), BW, Inches(1.7))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT
    box1.line.color.rgb = NAVY
    box1.line.width = Pt(1.2)

    add_text(slide, left=Inches(4.25), top=Inches(1.85), width=Inches(2.0), height=Inches(0.3),
             text="🌐 서비스 URL", size=11, bold=True, color=NAVY)
    add_text(slide, left=Inches(4.25), top=Inches(2.2), width=Inches(8.6), height=Inches(0.4),
             text="https://client-eight-ivory-60.vercel.app", size=14, bold=True, color=DARK)

    add_text(slide, left=Inches(4.25), top=Inches(2.75), width=Inches(2.0), height=Inches(0.3),
             text="📦 GitHub Repository", size=11, bold=True, color=NAVY)
    add_text(slide, left=Inches(4.25), top=Inches(3.05), width=Inches(8.6), height=Inches(0.4),
             text="https://github.com/Pionia5375/2026-1-web-programming", size=14, bold=True, color=DARK)

    # 박스 2: 핵심 기능
    add_text(slide, left=BX, top=Inches(3.7), width=BW, height=Inches(0.35),
             text="핵심 기능", size=14, bold=True, color=NAVY)
    add_bullet_block(
        slide, left=BX, top=Inches(4.05), width=BW, height=Inches(1.4),
        items=[
            "이메일/비밀번호 회원가입·로그인 (bcrypt + JWT 7일)",
            "할 일(Task) CRUD — 마감일, 우선순위, 검색·필터, 완료 토글",
            "일정(Event) CRUD — 시작/종료/장소",
            "통합 캘린더 (월/주/일/아젠다, 한국어) + 오늘 마감·오늘 일정 대시보드",
            "다크 모드, 마감 1시간 전 브라우저 알림(Notification API)",
        ],
        size=12,
    )

    # 박스 3: 기술 스택
    add_text(slide, left=BX, top=Inches(5.55), width=BW, height=Inches(0.35),
             text="기술 스택", size=14, bold=True, color=NAVY)
    add_bullet_block(
        slide, left=BX, top=Inches(5.9), width=BW, height=Inches(1.5),
        items=[
            "Frontend: React + Vite + Tailwind + react-big-calendar (Vercel 배포)",
            "Backend: Express + Prisma + JWT + bcrypt + zod (Render Docker 배포)",
            "Database: PostgreSQL (Neon)  /  Reverse Proxy: Nginx (gzip + SPA fallback)",
            "Web Storage: localStorage(JWT·테마) · sessionStorage(검색·뷰) · IndexedDB(오프라인 캐시)",
            "CI/CD: GitHub Actions (build·syntax·docker compose) + 자동 배포",
        ],
        size=12,
    )

    # 상단 우측: 강조 뱃지
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.6), Inches(0.4), Inches(1.55), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = AMBER
    badge.line.fill.background()
    add_text(slide, left=Inches(11.6), top=Inches(0.45), width=Inches(1.55), height=Inches(0.45),
             text="LIVE on Vercel + Render", size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    out = "/Users/ijeong/Desktop/웹 프로그래밍/2026-1-web-programming/docs/submission.pptx"
    pres.save(out)
    print(f"saved: {out}")


if __name__ == "__main__":
    main()
